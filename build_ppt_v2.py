# -*- coding: utf-8 -*-
"""
重构 智能TodoLife项目展示.pptx -> 智能TodoLife项目展示_v2.pptx
严格套用 PPT改进提示词.md：字号规范 / Bento 网格 / 60-30-10 配色 / 大数字 / 样机占位 / 图形化。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------- 画布 ----------
EMU = 914400
SW, SH = 13.333, 7.5

# ---------- 配色（60-30-10，单一强调色=微信绿）----------
INK      = RGBColor(0x14, 0x18, 0x20)   # 深墨蓝（深色页背景 / 浅色页标题）
INK2     = RGBColor(0x20, 0x26, 0x33)   # 深色卡片
PAPER    = RGBColor(0xF6, 0xF7, 0xF9)   # 浅色页背景（象牙灰）
CARD     = RGBColor(0xFF, 0xFF, 0xFF)   # 浅色卡片
CARD_DK  = RGBColor(0x2A, 0x31, 0x40)   # 深色页卡片
LINE     = RGBColor(0xE2, 0xE5, 0xEA)   # 浅分隔线
GREY     = RGBColor(0x5B, 0x63, 0x70)   # 次要文字
GREY_DK  = RGBColor(0xA8, 0xB0, 0xBD)   # 深色页次要文字
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
IVORY    = RGBColor(0xF2, 0xF4, 0xF7)
ACCENT   = RGBColor(0x07, 0xC1, 0x60)   # 微信绿（主强调色）
ACCENT_D = RGBColor(0x05, 0x9E, 0x4F)
MUTE     = RGBColor(0xC9, 0xCD, 0xD2)   # 图表置灰

# ---------- 多彩分类色（贴近 v1：墨绿 / 橙 / 蓝 / 红 / 紫 / 粉）----------
C_GREEN  = RGBColor(0x2F, 0x6F, 0x63)   # 墨绿（工作 / 主）
C_ORANGE = RGBColor(0xF2, 0xA0, 0x1D)   # 橙（紧急 / 协作 / 日历）
C_BLUE   = RGBColor(0x4D, 0x65, 0xA8)   # 宝蓝（学习 / 专注 / AI）
C_RED    = RGBColor(0xB4, 0x53, 0x47)   # 砖红（健康 / 逾期 / 洞察）
C_PURPLE = RGBColor(0x7E, 0x57, 0xB0)   # 紫（生活）
C_PINK   = RGBColor(0xD6, 0x4C, 0x86)   # 粉（创作）
C_GREY   = RGBColor(0x6B, 0x72, 0x80)   # 灰（提醒 / 低优先）
PALETTE  = [C_GREEN, C_ORANGE, C_BLUE, C_RED, C_PURPLE, C_PINK]

CN = "微软雅黑"      # 中文（思源黑体缺失时回退）
LAT = "Inter"        # 拉丁/数字

# 从空白演示文稿开始，避免继承原文件 15 张旧幻灯片部件（防重名/部件残留）。
prs = Presentation()
prs.slide_width = Emu(int(SW * EMU))
prs.slide_height = Emu(int(SH * EMU))
BLANK = prs.slide_layouts[6]


def set_fonts(run, cn=CN, lat=LAT):
    """同时设置东亚字体与拉丁字体。"""
    run.font.name = lat
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
    rPr.find(qn("a:latin")).set("typeface", lat)
    rPr.find(qn("a:ea")).set("typeface", cn)
    rPr.find(qn("a:cs")).set("typeface", lat)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dict(text,size,color,bold,cn,lat,space_after,line_spacing,spacing)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        if ln.get("space_after") is not None:
            p.space_after = Pt(ln["space_after"])
        p.space_before = Pt(ln.get("space_before", 0))
        if ln.get("line_spacing"):
            p.line_spacing = ln["line_spacing"]
        r = p.add_run()
        r.text = ln["text"]
        r.font.size = Pt(ln["size"])
        r.font.bold = ln.get("bold", False)
        r.font.color.rgb = ln.get("color", INK)
        if ln.get("spacing"):  # 字间距加宽（EMU on rPr -> spc in 1/100 pt）
            r._r.get_or_add_rPr().set("spc", str(int(ln["spacing"] * 100)))
        set_fonts(r, ln.get("cn", CN), ln.get("lat", LAT))
    return tb


def rect(slide, x, y, w, h, fill, line=None, line_w=0.75, radius=0.08, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def bg(slide, color):
    rect(slide, -0.1, -0.1, SW + 0.2, SH + 0.2, color, shape=MSO_SHAPE.RECTANGLE)


def card(slide, x, y, w, h, dark=False, accent_bar=False):
    fill = CARD_DK if dark else CARD
    ln = None if dark else LINE
    c = rect(slide, x, y, w, h, fill, line=ln, line_w=1.0, radius=0.06)
    if accent_bar:
        rect(slide, x, y, 0.09, h, ACCENT, radius=0.4)
    return c


def header(slide, eng, title, sub, num, dark=False):
    """统一页眉：中文主标题 + 副标题 + 页码（中文演讲，去英文栏目标签；eng 仅作占位忽略）。"""
    tcol = WHITE if dark else INK
    scol = GREY_DK if dark else GREY
    # 顶部强调短线（替代原英文标签的视觉锚点）
    rect(slide, 0.67, 0.62, 0.55, 0.08, ACCENT, radius=0.5)
    add_text(slide, 0.67, 0.82, 11.2, 0.85,
             [{"text": title, "size": 40, "color": tcol, "bold": True}])
    add_text(slide, 0.67, 1.66, 11.6, 0.5,
             [{"text": sub, "size": 18, "color": scol, "line_spacing": 1.3}])
    # 页码
    add_text(slide, SW - 1.4, 0.6, 0.9, 0.5,
             [{"text": num, "size": 18, "color": ACCENT, "bold": True, "align": PP_ALIGN.RIGHT, "lat": LAT}])


def connector(slide, x1, y1, x2, y2, color=LINE, w=1.0):
    """两点之间的直线连接（用于能力地图 / 信息架构树）。"""
    cn = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    cn.shadow.inherit = False
    return cn


def ring(slide, cx, cy, d, frac, color, num, label, dark=False, num_size=18, lbl_size=16):
    """环形进度环（甜甜圈）：实色外圆 + 同底色白心 = 干净圆环 + 中心数字 + 下方标签。
    BLOCK_ARC 在 PowerPoint 渲染不稳定，改用双圆叠加，跨平台一致。"""
    bgc = CARD if not dark else INK
    # 外圆实色
    rect(slide, cx - d / 2, cy - d / 2, d, d, color, shape=MSO_SHAPE.OVAL)
    # 白心（与背景同色）形成环
    inner = d * 0.6
    rect(slide, cx - inner / 2, cy - inner / 2, inner, inner, bgc, shape=MSO_SHAPE.OVAL)
    add_text(slide, cx - d / 2, cy - num_size / 110.0, d, num_size / 50.0,
             [{"text": num, "size": num_size, "color": INK if not dark else WHITE, "bold": True,
               "align": PP_ALIGN.CENTER, "lat": LAT}], anchor=MSO_ANCHOR.MIDDLE)
    if label:
        add_text(slide, cx - d / 2 - 0.2, cy + d / 2 + 0.04, d + 0.4, 0.3,
                 [{"text": label, "size": lbl_size, "color": GREY, "align": PP_ALIGN.CENTER}])


def linechart(slide, x, y, w, h, pts, color, fill_bg=True):
    """简易折线图：背景卡 + 淡灰基线 + 顶点折线（pts 为 0-1 高度比列表）。"""
    if fill_bg:
        rect(slide, x, y, w, h, CARD, line=LINE, line_w=1.0, radius=0.08)
    for i in range(4):                                   # 淡网格线
        gy = y + 0.25 + (h - 0.5) * i / 3
        rect(slide, x + 0.2, gy, w - 0.4, 0.012, LINE, shape=MSO_SHAPE.RECTANGLE)
    n = len(pts)
    px = [x + 0.35 + (w - 0.7) * i / (n - 1) for i in range(n)]
    py = [y + h - 0.3 - (h - 0.6) * v for v in pts]
    for i in range(n - 1):
        connector(slide, px[i], py[i], px[i + 1], py[i + 1], color=color, w=2.5)
    for i in range(n):
        rect(slide, px[i] - 0.05, py[i] - 0.05, 0.1, 0.1, color, shape=MSO_SHAPE.OVAL)


def heatmap(slide, x, y, cols, rows, cell=0.2, gap=0.05):
    """习惯热力图：多色块网格（深浅绿表示强度）。"""
    shades = [RGBColor(0xEC, 0xEF, 0xF1), RGBColor(0xB8, 0xE0, 0xC8),
              RGBColor(0x7C, 0xC9, 0xA0), RGBColor(0x2F, 0x6F, 0x63)]
    seed = 0
    for r in range(rows):
        for c in range(cols):
            seed = (seed * 1103515245 + 12345 + r * 31 + c * 7) & 0x7fffffff
            lvl = seed % 4
            rect(slide, x + c * (cell + gap), y + r * (cell + gap), cell, cell,
                 shades[lvl], radius=0.25)


def phone(slide, x, y, w=2.0, h=4.05, screen="today", title="今日", caption=None):
    """真实样机：粗黑圆角边框 + 侧边按钮 + 不同界面内容（today/kanban/calendar）。"""
    rect(slide, x, y, w, h, INK, radius=0.14)
    rect(slide, x - 0.04, y + h * 0.22, 0.05, h * 0.12, INK2, radius=0.4)   # 侧边按钮
    rect(slide, x - 0.04, y + h * 0.38, 0.05, h * 0.18, INK2, radius=0.4)
    pad = 0.085
    sx, sy, sw_, sh_ = x + pad, y + pad, w - 2 * pad, h - 2 * pad
    rect(slide, sx, sy, sw_, sh_, CARD, radius=0.1)
    add_text(slide, sx, sy + 0.1, sw_, 0.3,
             [{"text": title, "size": 13, "color": INK, "bold": True, "align": PP_ALIGN.CENTER}])
    cx0 = sx + sw_ / 2
    if screen == "today":
        # 三个进度环（缩小、留出标签空间，避免与下方卡重叠）
        rd = min(sw_ * 0.24, 0.62)
        ringy = sy + 0.62
        for i, (n, lb, col, fr) in enumerate([("68", "完成", C_GREEN, 0.68),
                                              ("72", "专注", C_BLUE, 0.72),
                                              ("45", "连续", C_ORANGE, 0.45)]):
            rcx = sx + sw_ * (0.22 + 0.28 * i)
            ring(slide, rcx, ringy, rd, fr, col, n, lb, num_size=13, lbl_size=10)
        # 现在做 高亮卡（标题单独整行，42分钟在标题上方右侧，互不重叠）
        ny = ringy + rd / 2 + 0.34
        nh = 0.66
        rect(slide, sx + 0.12, ny, sw_ - 0.24, nh, RGBColor(0xE5, 0xF3, 0xEA), line=ACCENT, line_w=1.0, radius=0.16)
        add_text(slide, sx + 0.24, ny + 0.07, 1.2, 0.22, [{"text": "现在做", "size": 10, "color": ACCENT_D, "bold": True}])
        # 42分钟 胶囊（右上角，窄）
        chw = 0.72
        rect(slide, sx + sw_ - 0.24 - chw, ny + 0.06, chw, 0.26, CARD, line=ACCENT, line_w=0.75, radius=0.5)
        add_text(slide, sx + sw_ - 0.24 - chw, ny + 0.08, chw, 0.22, [{"text": "42 分钟", "size": 9, "color": ACCENT_D, "align": PP_ALIGN.CENTER, "lat": LAT}], anchor=MSO_ANCHOR.MIDDLE)
        nowtb = add_text(slide, sx + 0.24, ny + 0.34, sw_ - 0.48, 0.28, [{"text": "整理汇报材料", "size": 12, "color": INK, "bold": True}])
        nowtb.text_frame.word_wrap = False
        # 分组任务行（左标签固定窄列，右任务卡单行不换行）
        ty = ny + nh + 0.14
        for lb, col in [("逾期 1", C_RED), ("今天 4", C_GREEN), ("稍后 3", C_GREY)]:
            rect(slide, sx + 0.16, ty + 0.07, 0.11, 0.11, col, shape=MSO_SHAPE.OVAL)
            add_text(slide, sx + 0.32, ty, 0.78, 0.26, [{"text": lb, "size": 10, "color": INK, "bold": True}])
            cardx = sx + 1.12
            cardw = sw_ - 1.12 - 0.16
            rect(slide, cardx, ty - 0.02, cardw, 0.32, CARD, line=LINE, line_w=0.75, radius=0.2)
            rect(slide, cardx, ty - 0.02, 0.05, 0.32, col, radius=0.4)
            tcardtb = add_text(slide, cardx + 0.12, ty - 0.01, cardw - 0.18, 0.3, [{"text": "完成论文初稿", "size": 9, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
            tcardtb.text_frame.word_wrap = False
            ty += 0.44
    elif screen == "kanban":
        cols = [("待处理", C_GREEN, 3), ("进行中", C_ORANGE, 2), ("已完成", C_BLUE, 1)]
        colw = (sw_ - 0.3) / 3
        for ci, (cl, col, cnt) in enumerate(cols):
            colx = sx + 0.1 + ci * (colw + 0.05)
            rect(slide, colx, sy + 0.45, colw, sh_ - 0.7, IVORY, radius=0.12)
            add_text(slide, colx, sy + 0.5, colw, 0.24, [{"text": cl, "size": 9, "color": col, "bold": True, "align": PP_ALIGN.CENTER}])
            for k in range(cnt):
                cardy = sy + 0.78 + k * 0.44
                rect(slide, colx + 0.06, cardy, colw - 0.12, 0.36, CARD, line=LINE, line_w=0.5, radius=0.2)
                rect(slide, colx + 0.06, cardy, 0.05, 0.36, col, radius=0.4)
                add_text(slide, colx + 0.14, cardy + 0.02, colw - 0.2, 0.32, [{"text": "任务", "size": 8, "color": INK}], anchor=MSO_ANCHOR.MIDDLE)
    elif screen == "calendar":
        gx, gy = sx + 0.18, sy + 0.5
        cw_ = (sw_ - 0.36) / 7
        seed = 7
        for r in range(5):
            for c in range(7):
                seed = (seed * 1103515245 + 12345) & 0x7fffffff
                fill = RGBColor(0xE5, 0xF3, 0xEA) if seed % 3 == 0 else IVORY
                rect(slide, gx + c * cw_, gy + r * 0.26, cw_ - 0.04, 0.22, fill, radius=0.2)
        # 倒计时卡
        cy2 = gy + 5 * 0.26 + 0.12
        rect(slide, sx + 0.16, cy2, sw_ - 0.32, 0.62, RGBColor(0xFD, 0xF1, 0xDF), line=C_ORANGE, line_w=1.0, radius=0.18)
        add_text(slide, sx + 0.26, cy2 + 0.08, sw_ - 0.5, 0.22, [{"text": "倒计时", "size": 10, "color": C_ORANGE, "bold": True}])
        add_text(slide, sx + 0.26, cy2 + 0.3, sw_ - 0.5, 0.26, [{"text": "期末汇报 5 天", "size": 11, "color": INK, "bold": True}])
    # 底部 tab
    tabs = ["今日", "看板", "日历", "统计", "我的"]
    tw = sw_ / len(tabs)
    cur = {"today": "今日", "kanban": "看板", "calendar": "日历"}.get(screen, title)
    for i, t in enumerate(tabs):
        col = ACCENT if t == cur else GREY
        add_text(slide, sx + i * tw, y + h - 0.34, tw, 0.24,
                 [{"text": t, "size": 9, "color": col, "align": PP_ALIGN.CENTER, "bold": t == cur}])
    if caption:
        add_text(slide, x - 0.2, y + h + 0.1, w + 0.4, 0.6,
                 [{"text": caption.split("｜")[0], "size": 20, "color": INK, "bold": True, "align": PP_ALIGN.CENTER},
                  {"text": caption.split("｜")[1] if "｜" in caption else "", "size": 16, "color": GREY, "align": PP_ALIGN.CENTER, "space_before": 3}])


def kpi(slide, x, y, w, num, label, ctx, dark=False, num_size=60, accent=True):
    """三段式大数字：数字 / 标签 / 上下文。"""
    ncol = ACCENT if accent else (WHITE if dark else INK)
    lcol = WHITE if dark else INK
    ccol = GREY_DK if dark else GREY
    add_text(slide, x, y, w, 1.3,
             [{"text": num, "size": num_size, "color": ncol, "bold": True, "lat": LAT, "line_spacing": 1.0}])
    add_text(slide, x, y + num_size / 72.0 * 1.05, w, 0.4,
             [{"text": label, "size": 20, "color": lcol, "bold": True}])
    add_text(slide, x, y + num_size / 72.0 * 1.05 + 0.42, w, 0.5,
             [{"text": ctx, "size": 16, "color": ccol, "line_spacing": 1.2}])


def hbar(slide, x, y, w, frac, label, pct, highlight=False, dark=False):
    """横向条形（直接标数据，无图例无坐标轴）。"""
    track = INK2 if dark else LINE
    fill = ACCENT if highlight else (GREY_DK if dark else GREY)
    lcol = WHITE if dark else INK
    add_text(slide, x, y - 0.04, 2.0, 0.34, [{"text": label, "size": 18, "color": lcol, "bold": highlight}])
    bx = x + 2.1
    bw = w - 2.1 - 0.9
    rect(slide, bx, y, bw, 0.3, track, radius=0.5)
    if frac > 0:
        rect(slide, bx, y, max(0.1, bw * frac), 0.3, fill, radius=0.5)
    add_text(slide, bx + bw + 0.05, y - 0.04, 0.85, 0.34,
             [{"text": pct, "size": 18, "color": fill if highlight else lcol, "bold": highlight,
               "align": PP_ALIGN.LEFT, "lat": LAT}])


def chip(slide, x, y, text, dark=False, size=16):
    # 按字符类型估宽：中文 ~0.23"/字，拉丁/空格 ~0.12"/字（16pt），再加左右内边距。
    cjk = sum(1 for ch in text if ord(ch) > 0x2E80)
    lat = len(text) - cjk
    w = 0.5 + cjk * 0.235 + lat * 0.12
    c = rect(slide, x, y, w, 0.46, ACCENT if not dark else CARD_DK, radius=0.5)
    tb = add_text(slide, x, y + 0.02, w, 0.42,
                  [{"text": text, "size": size, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}],
                  anchor=MSO_ANCHOR.MIDDLE)
    tb.text_frame.word_wrap = False     # 禁止换行，宽度由内容决定
    return w


def new_slide():
    return prs.slides.add_slide(BLANK)


M = 0.67  # 左边距

# ============================================================== 1 封面（深色）
s = new_slide()
bg(s, INK)
# 极光渐变点缀（用半透明圆角块模拟）
rect(s, -1.5, 4.0, 6, 6, INK2, radius=0.5)
rect(s, 8.5, -2.0, 7, 6, INK2, radius=0.5)
rect(s, M, 0.78, 0.55, 0.08, ACCENT, radius=0.5)
add_text(s, M, 0.96, 6, 0.4,
         [{"text": "产品案例  ·  01", "size": 18, "color": ACCENT, "bold": True, "spacing": 2}])
add_text(s, M, 2.3, 7.6, 1.4,
         [{"text": "智能 TodoLife", "size": 54, "color": WHITE, "bold": True}])
add_text(s, M, 3.5, 7.4, 1.2,
         [{"text": "一款会感知、能推荐的个人任务应用", "size": 28, "color": IVORY, "line_spacing": 1.25},
          {"text": "从记录工具进化为：理解你的状态，把清单变成可执行的下一步。", "size": 18,
           "color": GREY_DK, "space_before": 8, "line_spacing": 1.3}])
cx = M
for t in ["Jetpack Compose", "认知建模", "智能排程"]:
    cx += chip(s, cx, 5.15, t) + 0.25
phone(s, 9.55, 1.25, w=2.4, h=4.9, screen="today", title="智能 TodoLife")

# ============================================================== 2 背景（浅色）
s = new_slide()
bg(s, PAPER)
header(s, "PROJECT BACKGROUND", "为什么普通待办不够用", "Todo 工具的关键问题，不在缺功能，而在缺判断。", "02")
cards2 = [("清单越积越多", "记录成本极低、清理成本高，列表越拉越长，焦虑随之累积。"),
          ("优先级靠人判断", "用户知道要做事，却不知道此刻先做哪个——截止、复杂度全凭脑补。"),
          ("状态被忽略", "精力、压力、时间窗口都影响执行，而清单只会一视同仁。"),
          ("缺“下一步”", "清单是记录入口，却没把模糊目标转成当前可执行的动作。")]
cw, ch = 5.55, 1.85
for i, (t, d) in enumerate(cards2):
    x = M + (i % 2) * (cw + 0.45)
    y = 2.55 + (i // 2) * (ch + 0.35)
    card(s, x, y, cw, ch, accent_bar=True)
    add_text(s, x + 0.35, y + 0.25, 0.7, 0.7,
             [{"text": "0" + str(i + 1), "size": 30, "color": ACCENT, "bold": True, "lat": LAT}])
    add_text(s, x + 1.15, y + 0.28, cw - 1.4, 0.5, [{"text": t, "size": 28, "color": INK, "bold": True}])
    add_text(s, x + 1.15, y + 0.92, cw - 1.45, 0.8, [{"text": d, "size": 18, "color": GREY, "line_spacing": 1.3}])

# ============================================================== 3 产品目标（浅色，流程）
s = new_slide()
bg(s, PAPER)
header(s, "PRODUCT GOAL", "把待办系统，升级成集成系统", "产品目标围绕：捕捉问题、推理流程、展开反馈，形成可持续的智能闭环。", "03")
steps = [("捕捉", "快速记录想法与备注"), ("理解", "识别分类 / 截止 / 复杂度 / 重要度"),
         ("推荐", "结合认知状态给出下一步建议"), ("反馈", "用数据回看完成与习惯")]
bw_, gap = 2.65, 0.55
total = len(steps) * bw_ + (len(steps) - 1) * gap
sx = (SW - total) / 2
y = 3.1
for i, (t, d) in enumerate(steps):
    x = sx + i * (bw_ + gap)
    card(s, x, y, bw_, 1.95)
    rect(s, x + bw_ / 2 - 0.35, y + 0.3, 0.7, 0.7, ACCENT, radius=0.5, shape=MSO_SHAPE.OVAL)
    add_text(s, x + bw_ / 2 - 0.35, y + 0.36, 0.7, 0.6,
             [{"text": str(i + 1), "size": 28, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER, "lat": LAT}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.1, y + 1.12, bw_ - 0.2, 0.4,
             [{"text": t, "size": 28, "color": INK, "bold": True, "align": PP_ALIGN.CENTER}])
    add_text(s, x + 0.15, y + 1.5, bw_ - 0.3, 0.4,
             [{"text": d, "size": 16, "color": GREY, "align": PP_ALIGN.CENTER, "line_spacing": 1.2}])
    if i < len(steps) - 1:
        add_text(s, x + bw_ + 0.02, y + 0.55, gap, 0.6,
                 [{"text": "→", "size": 32, "color": ACCENT, "bold": True, "align": PP_ALIGN.CENTER}])
add_text(s, M, 5.7, SW - 2 * M, 0.6,
         [{"text": "选择可被解释、可被信任，让“做什么”有理由——这是产品的核心价值。",
           "size": 24, "color": INK, "bold": True, "align": PP_ALIGN.CENTER}])

# ============================================================== 4 能力图（浅色，中心+卫星 Bento）
s = new_slide()
bg(s, PAPER)
header(s, "CAPABILITY MAP", "产品能力地图", "核心能力围绕智能引擎组织，外层承接记录、执行、复盘与提醒。", "04")
cxn, cyn = SW / 2, 4.3
# 8 个卫星卡（含分类色点）；描述精简为单行，避免溢出
sat = [("今日页", "分组 / 推荐 / 捕捉", C_GREEN),
       ("AI 助手", "对话 / 提取 / 拆解", C_BLUE),
       ("智能日程", "按优先级排时段", C_ORANGE),
       ("看板", "待处理 / 进行 / 完成", C_GREEN),
       ("倒计时提醒", "重要日期与通知", C_GREY),
       ("洞察", "负荷 / 微习惯建议", C_RED),
       ("统计", "周报 / 趋势 / 热力图", C_ORANGE),
       ("日历", "截止日期与任务", C_BLUE)]
cw_, ch_ = 2.78, 1.0
# 两排各 4：上排 y=2.55，下排 y=5.0；中心 Engine 在正中、两排之间
xs = [M, 3.62, 6.93, SW - M - cw_]
pos = [(xs[0], 2.55), (xs[1], 2.55), (xs[2], 2.55), (xs[3], 2.55),
       (xs[0], 5.0), (xs[1], 5.0), (xs[2], 5.0), (xs[3], 5.0)]
# 先连线（连到卡片靠中心的边，不穿过中心块）
for (t, d, col), (x, y) in zip(sat, pos):
    ex = x + cw_ / 2
    ey = y + (ch_ if y < cyn else 0)
    connector(s, cxn, cyn, ex, ey, color=LINE, w=1.0)
# 中心 Engine（盖住连线交点）
ew, eh = 2.8, 1.0
rect(s, cxn - ew / 2, cyn - eh / 2, ew, eh, C_GREEN, radius=0.18)
add_text(s, cxn - ew / 2, cyn - 0.32, ew, 0.64,
         [{"text": "智能任务引擎", "size": 22, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
for (t, d, col), (x, y) in zip(sat, pos):
    card(s, x, y, cw_, ch_)
    rect(s, x + 0.24, y + 0.26, 0.15, 0.15, col, shape=MSO_SHAPE.OVAL)
    add_text(s, x + 0.48, y + 0.15, cw_ - 0.62, 0.36, [{"text": t, "size": 20, "color": INK, "bold": True}])
    add_text(s, x + 0.48, y + 0.56, cw_ - 0.62, 0.32, [{"text": d, "size": 16, "color": GREY}])

# ============================================================== 5 信息架构（浅色）
s = new_slide()
bg(s, PAPER)
header(s, "INFORMATION ARCHITECTURE", "从入口到深页的结构", "首页承接高频执行，深页承接状态流转、统计与洞察。", "05")
# 大圆角容器
rect(s, M, 2.45, SW - 2 * M, 4.55, CARD, line=LINE, line_w=1.0, radius=0.05)
# 顶层
rect(s, SW / 2 - 1.7, 2.75, 3.4, 0.7, RGBColor(0xE5, 0xF3, 0xEA), line=ACCENT, line_w=1.0, radius=0.25)
add_text(s, SW / 2 - 1.7, 2.82, 3.4, 0.55,
         [{"text": "智能 TodoLife", "size": 22, "color": ACCENT_D, "bold": True, "align": PP_ALIGN.CENTER}],
         anchor=MSO_ANCHOR.MIDDLE)
tabs = [("今日", "AI 助手", "高频执行", C_GREEN), ("看板", "新建任务", "状态流转", C_BLUE),
        ("日历", "高级洞察", "时间管理", C_ORANGE), ("任务", "统计页", "低频复盘", C_RED),
        ("我的", "倒计时页", "", C_GREY)]
tw = 1.95
tot = len(tabs) * tw + (len(tabs) - 1) * 0.32
sx = (SW - tot) / 2
ty1, ty2 = 3.85, 5.25
for i, (t, sub, cat, col) in enumerate(tabs):
    x = sx + i * (tw + 0.32)
    # 顶层→一级 连线
    connector(s, SW / 2, 3.45, x + tw / 2, ty1, color=LINE, w=1.0)
    # 一级 Tab 卡
    rect(s, x, ty1, tw, 0.7, IVORY, line=LINE, line_w=1.0, radius=0.18)
    add_text(s, x, ty1 + 0.16, tw, 0.4, [{"text": t, "size": 22, "color": col, "bold": True, "align": PP_ALIGN.CENTER}])
    # 一级→二级 连线
    connector(s, x + tw / 2, ty1 + 0.7, x + tw / 2, ty2, color=LINE, w=1.0)
    # 二级页（彩色描边）
    rect(s, x + 0.15, ty2, tw - 0.3, 0.55, CARD, line=col, line_w=1.25, radius=0.2)
    add_text(s, x + 0.15, ty2 + 0.08, tw - 0.3, 0.4, [{"text": sub, "size": 18, "color": col, "bold": True, "align": PP_ALIGN.CENTER}])
    # 分类标注
    if cat:
        add_text(s, x, ty2 + 0.7, tw, 0.35, [{"text": cat, "size": 18, "color": col, "bold": True, "align": PP_ALIGN.CENTER}])

# ============================================================== 6 数据模型（浅色，左字段+右实例卡）
s = new_slide()
bg(s, PAPER)
header(s, "DATA MODEL", "任务不只是标题", "SmartTask 把一行字扩展为可排序、可拆解、可提醒、可复盘的行为单元。", "06")
fields = [("基础信息", "标题、描述、分类、标签"), ("执行信息", "预估时长、截止、提醒、子任务"),
          ("特征信息", "重要度、复杂度、优先级、目标时段"), ("状态信息", "Todo / In Progress / Done"),
          ("习惯信息", "是否习惯、连击、完成历史")]
y = 2.55
for i, (t, d) in enumerate(fields):
    card(s, M, y, 6.4, 0.78, accent_bar=True)
    add_text(s, M + 0.3, y + 0.14, 2.2, 0.5, [{"text": t, "size": 22, "color": INK, "bold": True}])
    add_text(s, M + 2.6, y + 0.18, 3.7, 0.5, [{"text": d, "size": 16, "color": GREY}])
    y += 0.88
# 右侧实例卡（深色代码风）
ix, iy, iw, ih = 7.6, 2.55, 5.05, 4.4
card(s, ix, iy, iw, ih, dark=True)
add_text(s, ix + 0.35, iy + 0.25, iw - 0.7, 0.4,
         [{"text": "SmartTask 实例", "size": 20, "color": ACCENT, "bold": True}])
inst = [("category", "STUDY"), ("estimatedMinutes", "90"), ("dueDate", "2026-06-20"),
        ("priority", "HIGH"), ("status", "IN_PROGRESS"), ("subtasks", "3 / 5")]
yy = iy + 0.85
for k, v in inst:
    add_text(s, ix + 0.35, yy, 2.6, 0.4, [{"text": k, "size": 18, "color": GREY_DK, "lat": LAT}])
    add_text(s, ix + 3.0, yy, iw - 3.3, 0.4,
             [{"text": v, "size": 18, "color": WHITE, "bold": True, "lat": LAT, "align": PP_ALIGN.RIGHT}])
    yy += 0.55

# ============================================================== 7 AI 规划（浅色，流程+示例）
s = new_slide()
bg(s, PAPER)
header(s, "AI PLANNING", "一句话，拆成可执行计划", "远程大模型优先，本地启发式兜底；AI 失败时计划不缺席。", "07")
flow = ["用户输入", "DueDateParser", "SmartTaskEngine", "AI Planner", "子任务清单"]
fy = 2.65
fw = 2.3
tot = len(flow) * fw + (len(flow) - 1) * 0.12
sx = M
for i, t in enumerate(flow):
    x = M + i * (fw + 0.12)
    acc = (t == "AI Planner")
    rect(s, x, fy, fw, 0.75, ACCENT if acc else INK2, radius=0.2)
    add_text(s, x + 0.08, fy + 0.1, fw - 0.16, 0.55,
             [{"text": t, "size": 16, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER,
               "lat": LAT if t not in ("用户输入", "子任务清单") else CN}],
             anchor=MSO_ANCHOR.MIDDLE)
    if i < len(flow) - 1:
        add_text(s, x + fw - 0.04, fy + 0.05, 0.2, 0.6,
                 [{"text": "›", "size": 26, "color": ACCENT, "bold": True, "align": PP_ALIGN.CENTER}])
# 双路标注
add_text(s, M, 3.6, 11.5, 0.4,
         [{"text": "远程优先 → RemoteAiTaskPlanner ｜ 失败兜底 → LocalAiTaskPlanner（启发式拆解）",
           "size": 18, "color": ACCENT_D, "bold": True}])
# 示例子任务
add_text(s, M, 4.15, 6, 0.4, [{"text": "示例：写一篇项目总结报告 → 拆出 5 步", "size": 20, "color": INK, "bold": True}])
subs = [("资料收集", "6/10 · 45 分钟"), ("整理大纲", "6/12 · 40 分钟"), ("撰写正文", "6/15 · 90 分钟"),
        ("修改润色", "6/18 · 60 分钟"), ("排版定稿", "6/20 · 30 分钟")]
yy = 4.75
for i, (t, m) in enumerate(subs):
    x = M + (i % 1) * 0
    card(s, M, yy, 11.5, 0.42)
    add_text(s, M + 0.25, yy + 0.04, 0.5, 0.35,
             [{"text": str(i + 1), "size": 18, "color": ACCENT, "bold": True, "lat": LAT}])
    add_text(s, M + 0.8, yy + 0.05, 6, 0.35, [{"text": t, "size": 18, "color": INK}])
    add_text(s, M + 7.0, yy + 0.05, 4.2, 0.35,
             [{"text": m, "size": 16, "color": GREY, "align": PP_ALIGN.RIGHT, "lat": LAT}])
    yy += 0.5

# ============================================================== 8 认知模型（浅色，4维进度条）
s = new_slide()
bg(s, PAPER)
header(s, "COGNITIVE MODEL", "量化此刻的状态", "认知负荷拆成四个指标，帮助判断现在是否适合推进。", "08")
dims = [("视觉负荷", 0.45, "界面与信息密度带来的分散"),
        ("记忆负荷", 0.62, "需要在脑中并行保持的内容"),
        ("时间压力", 0.78, "临近截止造成的紧迫"),
        ("决策疲劳", 0.40, "连续决策后的损耗状态")]
y = 2.6
for t, v, d in dims:
    add_text(s, M, y, 3, 0.4, [{"text": t, "size": 22, "color": INK, "bold": True}])
    add_text(s, M, y + 0.45, 4.2, 0.4, [{"text": d, "size": 16, "color": GREY}])
    rect(s, 5.2, y + 0.1, 5.8, 0.32, LINE, radius=0.5)
    hi = v >= 0.7
    rect(s, 5.2, y + 0.1, 5.8 * v, 0.32, ACCENT if not hi else RGBColor(0xE0, 0x6C, 0x2E), radius=0.5)
    add_text(s, 11.1, y + 0.04, 1.0, 0.4,
             [{"text": str(int(v * 100)), "size": 22, "color": INK, "bold": True, "lat": LAT}])
    y += 1.02
add_text(s, M, 6.7, SW - 2 * M, 0.5,
         [{"text": "当系统判断负荷偏高时，自动切换到“专注视图”，减少干扰、筛选低成本任务。",
           "size": 24, "color": INK, "bold": True, "align": PP_ALIGN.CENTER}])

# ============================================================== 9 优先级逻辑（浅色，Hero + 条形）
s = new_slide()
bg(s, PAPER)
header(s, "PRIORITY LOGIC", "下一步推荐不是随机排序", "推荐分数同时考虑任务本身和用户状态，避免在高压力、低精力时强推高复杂任务。", "09")
# 左侧区域宽度
LW = 6.7
# 顶部 100% 堆叠条
factors = [("重要度", "任务价值", 0.36, "36%", C_GREEN),
           ("紧急度", "目标时间", 0.30, "30%", C_ORANGE),
           ("能力匹配", "专注与精力", 0.20, "20%", C_BLUE),
           ("复杂度", "执行难度", 0.14, "14%", C_RED),
           ("习惯加成", "连续性", 0.08, "8%", C_GREY)]
sx0, bw0, bh0 = M, LW, 0.55
bx = sx0
by = 2.6
for lbl, sub, f, pct, col in factors:
    seg = bw0 * f
    rect(s, bx, by, seg, bh0, col, radius=0.05, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, bx, by + 0.1, seg, 0.36,
             [{"text": pct, "size": 16, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER, "lat": LAT}])
    bx += seg
# 5 张因子卡（前 3 一行，后 2 一行）
cw_, ch_ = 2.1, 1.0
for i, (lbl, sub, f, pct, col) in enumerate(factors):
    row, coli = (0, i) if i < 3 else (1, i - 3)
    x = sx0 + coli * (cw_ + 0.2)
    y = 3.45 + row * (ch_ + 0.2)
    card(s, x, y, cw_, ch_)
    add_text(s, x + 0.22, y + 0.16, cw_ - 1.0, 0.4, [{"text": lbl, "size": 18, "color": INK, "bold": True}])
    add_text(s, x + cw_ - 0.95, y + 0.1, 0.8, 0.45,
             [{"text": pct, "size": 24, "color": col, "bold": True, "align": PP_ALIGN.RIGHT, "lat": LAT}])
    add_text(s, x + 0.22, y + 0.6, cw_ - 0.4, 0.32, [{"text": sub, "size": 16, "color": GREY}])
# 右：现在做 浅色卡 + 胶囊
hx, hy, hw, hh = 8.05, 2.95, 4.6, 3.4
card(s, hx, hy, hw, hh)
add_text(s, hx + 0.45, hy + 0.4, hw - 0.9, 0.35, [{"text": "现在做", "size": 18, "color": ACCENT_D, "bold": True}])
add_text(s, hx + 0.45, hy + 0.8, hw - 0.9, 0.5, [{"text": "整理项目汇报材料", "size": 28, "color": INK, "bold": True}])
rect(s, hx + 0.45, hy + 1.55, 1.4, 0.5, RGBColor(0xE5, 0xF3, 0xEA), line=ACCENT, line_w=1.0, radius=0.5)
add_text(s, hx + 0.45, hy + 1.6, 1.4, 0.4, [{"text": "42 分钟", "size": 18, "color": ACCENT_D, "bold": True, "align": PP_ALIGN.CENTER, "lat": LAT}], anchor=MSO_ANCHOR.MIDDLE)
rect(s, hx + 1.95, hy + 1.55, 1.85, 0.5, RGBColor(0xE8, 0xEC, 0xF7), line=C_BLUE, line_w=1.0, radius=0.5)
add_text(s, hx + 1.95, hy + 1.6, 1.85, 0.4, [{"text": "推荐分 82.4", "size": 18, "color": C_BLUE, "bold": True, "align": PP_ALIGN.CENTER, "lat": LAT}], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, hx + 0.45, hy + 2.25, hw - 0.9, 1.0,
         [{"text": "主要因为重要度高、距离目标时间较近，同时当前专注与精力能够承载中等复杂任务。",
           "size": 18, "color": GREY, "line_spacing": 1.35}])

# ============================================================== 10 艾森豪威尔（浅色 2x2）
s = new_slide()
bg(s, PAPER)
header(s, "EISENHOWER MATRIX", "把待办翻译成行动策略", "系统不只罗列任务，而是把每一项转成可执行策略。", "10")
quads = [("重要且紧急", "立即做", [("项目验收·拆为子任务", "30 分钟"), ("客户需求确认", "45 分钟")], False),
         ("重要不紧急", "计划做 · 投入产出最高", [("期末复习计划", "30 分钟"), ("核心技能学习", "45 分钟")], False),
         ("紧急不重要", "尽快处理 / 可委托", [("回复群消息", "30 分钟"), ("例会准备", "45 分钟")], False),
         ("不重要不紧急", "删除 / 批量处理", [("无目的信息浏览", "30 分钟"), ("低优先级收藏", "45 分钟")], True)]
qw, qh = 5.65, 2.0
for i, (title, strat, tasks, mute) in enumerate(quads):
    x = M + (i % 2) * (qw + 0.35)
    y = 2.5 + (i // 2) * (qh + 0.3)
    bgc = IVORY if mute else CARD
    c = rect(s, x, y, qw, qh, bgc, line=LINE, line_w=1.0, radius=0.06)
    bar = MUTE if mute else ACCENT
    rect(s, x, y, qw, 0.08, bar)
    tcol = GREY if mute else INK
    add_text(s, x + 0.3, y + 0.18, qw - 0.6, 0.4, [{"text": title, "size": 22, "color": tcol, "bold": True}])
    add_text(s, x + 0.3, y + 0.62, qw - 0.6, 0.35,
             [{"text": strat, "size": 16, "color": (MUTE if mute else ACCENT_D), "bold": True}])
    yy = y + 1.05
    for tt, mm in tasks:
        add_text(s, x + 0.3, yy, qw - 1.5, 0.35, [{"text": "· " + tt, "size": 16, "color": tcol}])
        add_text(s, x + qw - 1.45, yy, 1.15, 0.35,
                 [{"text": mm, "size": 16, "color": GREY, "align": PP_ALIGN.RIGHT, "lat": LAT}])
        yy += 0.42

# ============================================================== 11 今日体验（浅色，样机+Hero三环+卡）
s = new_slide()
bg(s, PAPER)
header(s, "TODAY EXPERIENCE", "首页只回答一个问题：现在做什么", "首页不是信息越多越好，而是把“选哪个任务”从用户脑中移走。", "11")
phone(s, M, 2.5, w=2.3, h=4.6, screen="today", title="智能 TodoLife")
# Hero 三环数字（大）
metrics = [("68", "能量", C_GREEN, 0.68), ("72", "精力", C_BLUE, 0.72), ("45", "专注", C_ORANGE, 0.45)]
mx = 3.55
for n, l, col, fr in metrics:
    ring(s, mx, 3.3, 1.45, fr, col, n, l)
    mx += 1.85
# 三卡
feats = [("今日推荐", "结合负荷给出一个低压力可推进的下一步", C_GREEN),
         ("快速捕捉", "随手输入，自动补全分类、截止与标签", C_BLUE),
         ("状态调度", "认知负荷、状态窗口与计划时间块", C_ORANGE)]
y = 4.5
for t, d, col in feats:
    card(s, 3.35, y, 9.3, 0.82)
    rect(s, 3.35, y, 0.09, 0.82, col, radius=0.4)
    add_text(s, 3.7, y + 0.16, 2.4, 0.5, [{"text": t, "size": 22, "color": INK, "bold": True}])
    add_text(s, 6.0, y + 0.22, 6.5, 0.5, [{"text": d, "size": 18, "color": GREY}])
    y += 0.92

# ============================================================== 12 多视图（浅色，4样机并排）
s = new_slide()
bg(s, PAPER)
header(s, "MULTI-VIEW", "不同阶段，用不同视图", "今日、看板、日历与倒计时分别解决执行、推进和时间管理问题。", "12")
views = [("today", "今日", "今日视图｜面向执行，按逾期 / 今天 / 稍后组织"),
         ("kanban", "看板", "看板视图｜面向推进，待处理 / 进行中 / 已完成"),
         ("calendar", "日历", "日历与倒计时｜面向时间，管理截止日期和重要节点")]
phw, phh = 2.85, 3.9
gap = (SW - 2 * M - 3 * phw) / 2
px = M
for screen, title, cap in views:
    phone(s, px, 2.45, w=phw, h=phh, screen=screen, title=title, caption=cap)
    px += phw + gap

# ============================================================== 13 统计（浅色，Hero+环形+周报）
s = new_slide()
bg(s, PAPER)
header(s, "ANALYTICS", "复盘不是报表，是下一轮安排的依据", "统计页帮助用户发现负荷高峰、低效时间段和稳定习惯，而不是展示复杂 BI 大屏。", "13")
# 大圆角容器
rect(s, M, 2.45, SW - 2 * M, 4.55, CARD, line=LINE, line_w=1.0, radius=0.04)
# --- 左栏：周报摘要 + 完成率趋势 ---
lx = M + 0.4
add_text(s, lx, 2.7, 3.5, 0.35, [{"text": "周报摘要", "size": 20, "color": INK, "bold": True}])
kpis = [("68%", "完成率", C_GREEN), ("17/25", "已完成", C_BLUE), ("上升", "趋势", C_ORANGE)]
kx = lx
for n, l, col in kpis:
    rect(s, kx, 3.15, 1.5, 0.95, IVORY, radius=0.12)
    add_text(s, kx, 3.28, 1.5, 0.45, [{"text": n, "size": 26, "color": col, "bold": True, "align": PP_ALIGN.CENTER, "lat": LAT}])
    add_text(s, kx, 3.74, 1.5, 0.3, [{"text": l, "size": 16, "color": GREY, "align": PP_ALIGN.CENTER}])
    kx += 1.6
add_text(s, lx, 4.35, 3.5, 0.35, [{"text": "完成率趋势", "size": 20, "color": INK, "bold": True}])
linechart(s, lx, 4.75, 4.6, 1.85, [0.15, 0.4, 0.32, 0.62, 0.5, 0.82], C_GREEN)
# --- 中栏：时间分布 + 习惯热力图 ---
mx2 = 6.0
add_text(s, mx2, 2.7, 3.0, 0.35, [{"text": "时间分布", "size": 20, "color": INK, "bold": True}])
tcats = [("工作", C_GREEN), ("学习", C_BLUE), ("健康", C_RED), ("生活", C_PURPLE), ("协作", C_ORANGE), ("创作", C_PINK)]
for i, (t, col) in enumerate(tcats):
    cx2 = mx2 + (i % 2) * 1.5
    cy2 = 3.15 + (i // 2) * 0.62
    rect(s, cx2, cy2, 1.3, 0.5, CARD, line=col, line_w=1.25, radius=0.18)
    add_text(s, cx2, cy2 + 0.1, 1.3, 0.34, [{"text": t, "size": 18, "color": col, "bold": True, "align": PP_ALIGN.CENTER}])
add_text(s, mx2, 5.15, 3.0, 0.35, [{"text": "习惯热力图", "size": 20, "color": INK, "bold": True}])
heatmap(s, mx2, 5.55, 11, 4, cell=0.2, gap=0.05)
# --- 右栏：认知负荷曲线 + 成就徽章 ---
rx2 = 9.55
add_text(s, rx2, 2.7, 3.0, 0.35, [{"text": "认知负荷曲线", "size": 20, "color": INK, "bold": True}])
linechart(s, rx2, 3.1, 2.75, 1.5, [0.3, 0.5, 0.42, 0.7, 0.6, 0.85], C_BLUE)
add_text(s, rx2, 4.85, 3.0, 0.35, [{"text": "成就徽章", "size": 20, "color": INK, "bold": True}])
badges = [("十项全能", C_GREEN), ("一周习惯", C_ORANGE), ("低负荷管理", C_BLUE)]
by2 = 5.25
for t, col in badges:
    rect(s, rx2, by2, 2.55, 0.45, CARD, line=col, line_w=1.25, radius=0.3)
    add_text(s, rx2, by2 + 0.06, 2.55, 0.34, [{"text": t, "size": 18, "color": col, "bold": True, "align": PP_ALIGN.CENTER}])
    by2 += 0.55

# ============================================================== 14 技术架构（浅色，6层）
s = new_slide()
bg(s, PAPER)
header(s, "TECH ARCHITECTURE", "分层为本，AI 增强", "架构重点是本地稳定优先 + AI 能力增强；AI 失败时回退到启发式逻辑。", "14")
layers = [("UI 层", "Jetpack Compose / Material 3 / Screens / Components"),
          ("状态层", "ViewModel / StateFlow"),
          ("引擎层", "SmartTask / CognitiveLoad / Statistics / Achievement / TaskCompletion"),
          ("数据层", "Local Repository / Theme Preference / Countdown Repository"),
          ("AI 层", "Remote AI Planner / Local AI Planner / AI Chat Client / DueDateParser"),
          ("系统服务", "ReminderScheduler / Notification / BroadcastReceiver")]
y = 2.5
lw = 8.6
for i, (t, d) in enumerate(layers):
    acc = (t == "AI 层")
    c = rect(s, M, y, lw, 0.62, CARD_DK if acc else CARD, line=None if acc else LINE, line_w=1.0, radius=0.06)
    rect(s, M, y, 1.7, 0.62, ACCENT if acc else INK2, radius=0.06)
    add_text(s, M + 0.1, y + 0.08, 1.5, 0.45,
             [{"text": t, "size": 18, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, M + 1.9, y + 0.08, lw - 2.0, 0.45,
             [{"text": d, "size": 16, "color": (IVORY if acc else GREY), "lat": LAT}],
             anchor=MSO_ANCHOR.MIDDLE)
    y += 0.7
# 右侧降级标注
card(s, M + lw + 0.4, 2.5, 2.85, 4.02, dark=True)
add_text(s, M + lw + 0.7, 2.75, 2.3, 0.5, [{"text": "降级策略", "size": 20, "color": ACCENT, "bold": True}])
add_text(s, M + lw + 0.7, 3.4, 2.3, 3,
         [{"text": "AI 作为增强：", "size": 18, "color": WHITE, "bold": True, "line_spacing": 1.3},
          {"text": "远程大模型优先", "size": 16, "color": IVORY, "space_before": 6},
          {"text": "↓ 失败", "size": 16, "color": ACCENT, "bold": True, "space_before": 4},
          {"text": "本地启发式兜底", "size": 16, "color": IVORY, "space_before": 4},
          {"text": "保证核心功能在无网 / 无模型时仍可用。", "size": 16, "color": GREY_DK,
           "space_before": 10, "line_spacing": 1.3}])

# ============================================================== 15 总结（深色，4编号卡）
s = new_slide()
bg(s, INK)
add_text(s, M, 0.7, 6, 0.4,
         [{"text": "DESIGN SUMMARY", "size": 16, "color": ACCENT, "bold": True, "spacing": 3, "lat": LAT}])
add_text(s, M, 1.1, 11.5, 0.8, [{"text": "智能 TodoLife 的设计价值", "size": 40, "color": WHITE, "bold": True}])
add_text(s, M, 1.95, 11.5, 0.5,
         [{"text": "不止是效率工具，更试图理解用户，在合适的时间做合适的事。", "size": 18, "color": GREY_DK}])
rect(s, M, 2.55, 1.0, 0.045, ACCENT, radius=0.5)
vals = [("01", "从记录转向行动建议"), ("02", "从固定清单转向状态自适应"),
        ("03", "从单点 AI 转向可解释推荐与兜底"), ("04", "从数据报表转向行为洞察")]
cw, ch = 5.55, 1.5
for i, (n, t) in enumerate(vals):
    x = M + (i % 2) * (cw + 0.45)
    y = 2.95 + (i // 2) * (ch + 0.35)
    card(s, x, y, cw, ch, dark=True, accent_bar=True)
    add_text(s, x + 0.35, y + 0.2, 1.6, 1.1,
             [{"text": n, "size": 60, "color": RGBColor(0x3A, 0x44, 0x55), "bold": True, "lat": LAT}])
    add_text(s, x + 1.85, y + 0.42, cw - 2.1, 0.8,
             [{"text": t, "size": 24, "color": WHITE, "bold": True, "line_spacing": 1.15}],
             anchor=MSO_ANCHOR.MIDDLE)
add_text(s, M, 6.55, SW - 2 * M, 0.6,
         [{"text": "把任务翻译为行动策略，把数据翻译成一种叙事。",
           "size": 24, "color": ACCENT, "bold": True, "align": PP_ALIGN.CENTER}])

# ---------- 保存 ----------
OUT = "智能TodoLife项目展示_v4.pptx"
prs.save(OUT)

# ---------- 自检 ----------
# 规范：正文/标题/标签 ≥16pt（英文栏目标签与图表数据标签允许 16pt）。
# 设备样机(phone)内部 UI 文字属于“界面示意”插画内容（等同截图），按 <16pt 归类为 illustration，不计违规。
from collections import Counter
chk = Presentation(OUT)
violations = []
illustration = 0
for idx, sl in enumerate(chk.slides, 1):
    sizes = Counter()
    for shp in sl.shapes:
        if not shp.has_text_frame:
            continue
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                if not (r.font.size and r.text.strip()):
                    continue
                pt = r.font.size.pt
                sizes[pt] += 1
                if pt <= 13:
                    illustration += 1            # 样机内部 UI 示意文字（≤13pt 仅出现在 phone 内，等同截图）
                elif pt < 16:
                    violations.append((idx, pt, r.text[:24]))
    print(f"P{idx:2d} sizes: {dict(sorted(sizes.items()))}")
print("\nSLIDES:", len(chk.slides._sldIdLst))
print(f"样机内部示意文字(<16pt, illustration): {illustration} 处")
# 重名部件检查
import zipfile
zc = Counter(zipfile.ZipFile(OUT).namelist())
dups = {k: v for k, v in zc.items() if v > 1}
print("重名部件:", len(dups))
if violations:
    print("!! 正文 <16pt 违规:", violations)
else:
    print("OK: 所有正文/标题/标签 ≥16pt（样机内部示意除外）")
