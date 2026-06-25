from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from math import cos, sin, pi


OUT = "智能TodoLife项目展示.pptx"

W, H = 13.333, 7.5

GREEN = RGBColor(0x2F, 0x6F, 0x63)
WECHAT = RGBColor(0x07, 0xC1, 0x60)
BLUE = RGBColor(0x4D, 0x65, 0xA8)
RED = RGBColor(0xB4, 0x53, 0x47)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
TEXT = RGBColor(0x1F, 0x29, 0x33)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xF8)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD9, 0xE2, 0xDE)
PALE_GREEN = RGBColor(0xE9, 0xF5, 0xF0)
PALE_BLUE = RGBColor(0xEC, 0xF1, 0xFA)
PALE_RED = RGBColor(0xFA, 0xEE, 0xEC)
PALE_ORANGE = RGBColor(0xFE, 0xF4, 0xE4)
PALE_GRAY = RGBColor(0xF1, 0xF2, 0xF3)

FONT = "Microsoft YaHei"


def inch(v):
    return Inches(v)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color=LINE, width=1, transparency=0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, font=FONT, line_spacing=1.05):
    box = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_label(slide, text, x, y, w, h, color=GREEN, fill=PALE_GREEN, size=10):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    set_fill(shape, fill)
    set_line(shape, color, 0.7, 45)
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return shape


def add_panel(slide, x, y, w, h, fill=PANEL, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
    set_fill(shape, fill)
    set_line(shape, line, 0.8, 10)
    return shape


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_text(slide, section, 0.66, 0.36, 2.0, 0.25, 8.5, True, GREEN)
    add_text(slide, title, 0.66, 0.62, 8.8, 0.48, 24, True, TEXT)
    if subtitle:
        add_text(slide, subtitle, 0.66, 1.13, 8.9, 0.32, 10.5, False, MUTED)


def add_page_num(slide, idx):
    add_text(slide, f"{idx:02d}", 12.45, 7.06, 0.34, 0.16, 8, False, RGBColor(0xA7, 0xAD, 0xB2), PP_ALIGN.RIGHT)


def add_footer_line(slide):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(0.66), inch(6.98), inch(12.1), inch(6.98))
    ln.line.color.rgb = RGBColor(0xEA, 0xEE, 0xEC)
    ln.line.width = Pt(0.6)


def add_slide(prs, idx, title, subtitle=None, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inch(W), inch(H))
    set_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    add_title(slide, title, subtitle, section)
    add_footer_line(slide)
    add_page_num(slide, idx)
    return slide


def add_bullet(slide, title, body, x, y, w, h, accent=GREEN, fill=PANEL):
    add_panel(slide, x, y, w, h, fill=fill)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x + 0.18), inch(y + 0.22), inch(0.13), inch(0.13))
    set_fill(dot, accent)
    dot.line.fill.background()
    add_text(slide, title, x + 0.38, y + 0.16, w - 0.56, 0.25, 13, True, TEXT)
    add_text(slide, body, x + 0.38, y + 0.52, w - 0.56, h - 0.62, 9.8, False, MUTED)


def add_arrow(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def add_phone(slide, x, y, w, h, screen_title="今日", mode="today"):
    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    set_fill(phone, RGBColor(0x17, 0x1B, 0x1F))
    phone.line.fill.background()
    inner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x + 0.08), inch(y + 0.12), inch(w - 0.16), inch(h - 0.24))
    set_fill(inner, RGBColor(0xFB, 0xFC, 0xFB))
    inner.line.fill.background()
    add_text(slide, screen_title, x + 0.22, y + 0.28, w - 0.44, 0.25, 9, True, TEXT, PP_ALIGN.CENTER)

    if mode == "today":
        add_rings(slide, x + 0.25, y + 0.68, w - 0.5, 0.9, small=True)
        card = add_panel(slide, x + 0.22, y + 1.72, w - 0.44, 0.58, fill=PALE_GREEN, line=RGBColor(0xC7, 0xE6, 0xDA))
        add_text(slide, "现在做", x + 0.38, y + 1.82, 0.7, 0.13, 6.8, True, GREEN)
        add_text(slide, "整理项目汇报材料", x + 0.38, y + 2.01, w - 0.95, 0.17, 7.5, True, TEXT)
        add_label(slide, "42 分钟", x + w - 1.05, y + 1.93, 0.62, 0.2, GREEN, RGBColor(0xD8, 0xEF, 0xE5), 6)
        for i, (t, c) in enumerate([("逾期 1", RED), ("今天 4", GREEN), ("稍后 3", GRAY)]):
            yy = y + 2.55 + i * 0.48
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x + 0.28), inch(yy + 0.08), inch(0.08), inch(0.08))
            set_fill(dot, c)
            dot.line.fill.background()
            add_text(slide, t, x + 0.42, yy, 0.8, 0.15, 6.8, True, TEXT)
            add_task_bar(slide, x + 1.04, yy - 0.01, w - 1.38, 0.23, "完成期末论文初稿", c, 6.2)
        for i, t in enumerate(["今日", "看板", "日历", "任务", "我的"]):
            add_text(slide, t, x + 0.22 + i * ((w - 0.44) / 5), y + h - 0.48, (w - 0.44) / 5, 0.15, 5.6, i == 0, GREEN if i == 0 else MUTED, PP_ALIGN.CENTER)
    elif mode == "board":
        cols = [("待处理", GREEN), ("进行中", ORANGE), ("已完成", BLUE)]
        for i, (name, col) in enumerate(cols):
            xx = x + 0.22 + i * ((w - 0.52) / 3)
            cw = (w - 0.7) / 3
            add_panel(slide, xx, y + 0.72, cw, h - 1.45, fill=RGBColor(0xF5, 0xF7, 0xF6), line=RGBColor(0xE2, 0xE8, 0xE5))
            add_text(slide, name, xx + 0.08, y + 0.86, cw - 0.16, 0.14, 5.8, True, col, PP_ALIGN.CENTER)
            for j in range(3 - i):
                add_task_bar(slide, xx + 0.08, y + 1.16 + j * 0.4, cw - 0.16, 0.27, "任务", col, 5.2)
    elif mode == "calendar":
        for r in range(5):
            for c in range(7):
                xx, yy = x + 0.25 + c * 0.3, y + 0.8 + r * 0.28
                sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(xx), inch(yy), inch(0.22), inch(0.2))
                set_fill(sq, PALE_GREEN if (r + c) % 4 == 0 else RGBColor(0xF0, 0xF3, 0xF1))
                sq.line.fill.background()
        add_panel(slide, x + 0.25, y + 2.55, w - 0.5, 0.75, fill=PALE_ORANGE, line=RGBColor(0xF2, 0xD2, 0xA8))
        add_text(slide, "倒计时", x + 0.42, y + 2.7, 0.8, 0.16, 7, True, ORANGE)
        add_text(slide, "期末汇报 5 天", x + 0.42, y + 2.95, w - 0.84, 0.18, 7.8, True, TEXT)
    elif mode == "stats":
        add_text(slide, "周报摘要", x + 0.25, y + 0.7, w - 0.5, 0.18, 7.5, True, TEXT)
        for i, (v, lab) in enumerate([("68%", "完成率"), ("17/25", "已完成"), ("上升", "趋势")]):
            add_panel(slide, x + 0.25 + i * 0.74, y + 1.0, 0.62, 0.46, fill=PALE_GRAY, line=RGBColor(0xE4, 0xE8, 0xE6))
            add_text(slide, v, x + 0.29 + i * 0.74, y + 1.08, 0.54, 0.13, 7.2, True, GREEN, PP_ALIGN.CENTER)
            add_text(slide, lab, x + 0.29 + i * 0.74, y + 1.25, 0.54, 0.1, 4.9, False, MUTED, PP_ALIGN.CENTER)
        add_mini_line_chart(slide, x + 0.28, y + 1.72, w - 0.56, 0.9, GREEN)
        add_heatmap(slide, x + 0.32, y + 2.9, 8, 4, 0.12)
    return phone


def add_task_bar(slide, x, y, w, h, text, color=GREEN, size=7):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    set_fill(bar, RGBColor(0xFF, 0xFF, 0xFF))
    set_line(bar, RGBColor(0xE5, 0xEA, 0xE8), 0.5)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(0.035), inch(h))
    set_fill(strip, color)
    strip.line.fill.background()
    add_text(slide, text, x + 0.08, y + h * 0.25, w - 0.16, h * 0.5, size, False, TEXT)


def add_rings(slide, x, y, w, h, small=False):
    colors = [GREEN, BLUE, ORANGE]
    vals = [0.68, 0.72, 0.45]
    labels = ["完成", "专注", "连续"]
    r = min(h * 0.58, w / 5)
    gap = w / 3
    for i in range(3):
        cx = x + gap * i + gap / 2 - r / 2
        cy = y + 0.03
        base = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(cx), inch(cy), inch(r), inch(r))
        set_fill(base, RGBColor(0xF1, 0xF4, 0xF2))
        base.line.fill.background()
        arc = slide.shapes.add_shape(MSO_SHAPE.ARC, inch(cx + 0.02), inch(cy + 0.02), inch(r - 0.04), inch(r - 0.04))
        arc.adjustments[0] = -90
        arc.adjustments[1] = -90 + int(360 * vals[i])
        arc.line.color.rgb = colors[i]
        arc.line.width = Pt(4 if not small else 2.4)
        arc.fill.background()
        add_text(slide, f"{int(vals[i] * 100)}", cx + 0.12, cy + r * 0.34, r - 0.24, 0.15, 8 if small else 13, True, TEXT, PP_ALIGN.CENTER)
        add_text(slide, labels[i], cx, cy + r + 0.06, r, 0.13, 5.5 if small else 8, False, MUTED, PP_ALIGN.CENTER)


def add_mini_line_chart(slide, x, y, w, h, color=GREEN):
    add_panel(slide, x, y, w, h, fill=RGBColor(0xFF, 0xFF, 0xFF), line=RGBColor(0xE4, 0xE8, 0xE6))
    pts = [(0.08, 0.72), (0.2, 0.55), (0.34, 0.62), (0.5, 0.35), (0.68, 0.44), (0.84, 0.24)]
    for i in range(4):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x + 0.12), inch(y + h * (0.2 + i * 0.18)), inch(x + w - 0.12), inch(y + h * (0.2 + i * 0.18)))
        ln.line.color.rgb = RGBColor(0xEC, 0xEF, 0xEE)
        ln.line.width = Pt(0.4)
    for i in range(len(pts) - 1):
        x1, y1 = x + pts[i][0] * w, y + pts[i][1] * h
        x2, y2 = x + pts[i + 1][0] * w, y + pts[i + 1][1] * h
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
        ln.line.color.rgb = color
        ln.line.width = Pt(1.7)
    for px, py in pts:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x + px * w - 0.025), inch(y + py * h - 0.025), inch(0.05), inch(0.05))
        set_fill(dot, color)
        dot.line.fill.background()


def add_heatmap(slide, x, y, cols, rows, cell):
    vals = [0, 1, 2, 3, 1, 0, 2, 3, 0, 1, 1, 2, 3, 0, 1, 2, 0, 0, 3, 2, 1, 1, 0, 3, 2, 1, 0, 2, 3, 1, 0, 1]
    palette = [RGBColor(0xEA, 0xEE, 0xEC), RGBColor(0xC9, 0xE6, 0xDC), RGBColor(0x8F, 0xCC, 0xB7), GREEN]
    k = 0
    for r in range(rows):
        for c in range(cols):
            sq = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x + c * (cell + 0.035)), inch(y + r * (cell + 0.035)), inch(cell), inch(cell))
            set_fill(sq, palette[vals[k % len(vals)]])
            sq.line.fill.background()
            k += 1


def add_factor_card(slide, x, y, label, percent, color, note):
    add_panel(slide, x, y, 1.78, 0.78, fill=PANEL)
    add_text(slide, label, x + 0.15, y + 0.12, 1.0, 0.18, 9, True, TEXT)
    add_text(slide, percent, x + 1.13, y + 0.11, 0.45, 0.18, 13, True, color, PP_ALIGN.RIGHT)
    add_text(slide, note, x + 0.15, y + 0.41, 1.45, 0.2, 7.3, False, MUTED)


def build_ppt():
    prs = Presentation()
    prs.slide_width = inch(W)
    prs.slide_height = inch(H)

    # 01
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, inch(W), inch(H))
    set_fill(bg, LIGHT_BG)
    bg.line.fill.background()
    add_text(s, "智能 TodoLife", 0.78, 1.16, 5.3, 0.72, 32, True, TEXT)
    add_text(s, "认知负荷驱动的个人任务管理应用", 0.82, 1.98, 4.8, 0.35, 14, False, GREEN)
    add_text(s, "从记录任务升级为理解任务、安排任务、执行任务、复盘任务。产品重点不是创建更多事项，而是降低选择成本，把当下最该做的事推到用户面前。", 0.83, 2.56, 4.75, 0.9, 12, False, MUTED)
    add_label(s, "Jetpack Compose", 0.83, 3.82, 1.55, 0.34, GREEN, PALE_GREEN, 9)
    add_label(s, "认知负荷", 2.55, 3.82, 1.22, 0.34, BLUE, PALE_BLUE, 9)
    add_label(s, "智能排程", 3.95, 3.82, 1.22, 0.34, ORANGE, PALE_ORANGE, 9)
    add_phone(s, 8.05, 0.78, 2.65, 5.7, "智能 TodoLife", "today")
    add_panel(s, 7.68, 1.16, 0.1, 4.95, fill=GREEN, line=GREEN)
    add_text(s, "PRODUCT CASE", 0.84, 6.62, 1.6, 0.2, 8, True, RGBColor(0x9A, 0xA3, 0xA0))
    add_footer_line(s)
    add_page_num(s, 1)

    # 02
    s = add_slide(prs, 2, "为什么普通待办不够用", "Todo 工具的关键问题，不是缺少任务入口，而是缺少“下一步决策”。", "PROJECT BACKGROUND")
    items = [
        ("任务越记越多", "记录成本很低，但整理成本持续升高。用户打开清单后，面对的是一个越来越长的任务池。", GREEN, PALE_GREEN),
        ("优先级难判断", "用户知道“要做”，但常常不知道“先做哪个”。重要度、截止时间、任务复杂度混在一起。", BLUE, PALE_BLUE),
        ("状态被忽略", "精力、压力、环境和时间窗口会影响执行质量，但多数待办工具只处理任务本身。", ORANGE, PALE_ORANGE),
    ]
    for i, (t, b, c, f) in enumerate(items):
        add_bullet(s, t, b, 0.82 + i * 4.05, 2.0, 3.42, 2.42, c, f)
    add_panel(s, 1.24, 5.32, 10.85, 0.82, fill=PANEL)
    add_text(s, "设计判断", 1.55, 5.55, 1.0, 0.2, 10, True, GREEN)
    add_text(s, "待办清单的下一步不是更快录入，而是帮助用户把任务池转译成当前可执行的行动。", 2.55, 5.53, 8.6, 0.25, 13, True, TEXT)

    # 03
    s = add_slide(prs, 3, "把待办系统做成决策系统", "产品目标围绕“捕捉、理解、推荐、复盘”展开，形成可持续的任务管理闭环。", "PRODUCT GOAL")
    kws = [
        ("捕捉", "快速记录任务与备注", GREEN),
        ("理解", "识别分类、截止日期、复杂度、重要度", BLUE),
        ("推荐", "结合认知状态计算下一步任务", ORANGE),
        ("复盘", "用数据回看完成率、习惯、认知负荷趋势", RED),
    ]
    for i, (k, b, c) in enumerate(kws):
        add_bullet(s, k, b, 0.82, 1.72 + i * 1.0, 3.8, 0.72, c, PANEL)
    cx, cy, r = 8.8, 3.6, 1.65
    steps = [("输入任务", GREEN), ("智能分析", BLUE), ("推荐行动", ORANGE), ("执行反馈", RED), ("数据洞察", BLUE), ("优化安排", GREEN)]
    coords = []
    for i, (label, col) in enumerate(steps):
        a = -pi / 2 + i * 2 * pi / len(steps)
        x, y = cx + cos(a) * r, cy + sin(a) * r
        coords.append((x, y))
        node = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x - 0.58), inch(y - 0.22), inch(1.16), inch(0.44))
        set_fill(node, PANEL)
        set_line(node, col, 1)
        add_text(s, label, x - 0.48, y - 0.075, 0.96, 0.12, 7.8, True, col, PP_ALIGN.CENTER)
    for i in range(len(coords)):
        x1, y1 = coords[i]
        x2, y2 = coords[(i + 1) % len(coords)]
        add_arrow(s, x1 + (x2 - x1) * 0.26, y1 + (y2 - y1) * 0.26, x1 + (x2 - x1) * 0.74, y1 + (y2 - y1) * 0.74, RGBColor(0xB9, 0xC6, 0xC0), 1)
    add_panel(s, 7.55, 3.18, 2.5, 0.72, fill=PALE_GREEN, line=RGBColor(0xC6, 0xE4, 0xD9))
    add_text(s, "降低选择成本", 7.82, 3.32, 1.95, 0.18, 13, True, GREEN, PP_ALIGN.CENTER)

    # 04
    s = add_slide(prs, 4, "产品能力地图", "核心能力围绕 Smart Todo Engine 组织，外层承接用户的记录、执行、复盘和提醒。", "CAPABILITY MAP")
    center = add_panel(s, 5.25, 2.68, 2.75, 0.82, fill=GREEN, line=GREEN)
    add_text(s, "Smart Todo Engine", 5.52, 2.92, 2.2, 0.18, 13, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)
    mods = [
        ("今日页", "任务分组 / 当前推荐 / 快速捕捉", GREEN),
        ("AI 助手", "对话 / 任务提取 / 目标拆解", BLUE),
        ("智能日程", "按优先级和认知负荷安排时间段", ORANGE),
        ("看板", "待处理 / 进行中 / 已完成", GREEN),
        ("日历", "截止日期与计划任务", BLUE),
        ("统计", "周报 / 趋势 / 热力图", ORANGE),
        ("洞察", "认知负荷 / 微习惯建议", RED),
        ("倒计时提醒", "重要日期与通知", GRAY),
    ]
    positions = [(1.0, 1.55), (4.08, 1.25), (7.18, 1.25), (10.1, 1.55),
                 (10.1, 4.45), (7.18, 4.75), (4.08, 4.75), (1.0, 4.45)]
    for (name, desc, col), (x, y) in zip(mods, positions):
        add_panel(s, x, y, 2.28, 0.82, fill=PANEL)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, inch(x + 0.16), inch(y + 0.2), inch(0.12), inch(0.12))
        set_fill(dot, col)
        dot.line.fill.background()
        add_text(s, name, x + 0.36, y + 0.13, 1.7, 0.16, 9.6, True, TEXT)
        add_text(s, desc, x + 0.36, y + 0.38, 1.74, 0.22, 6.8, False, MUTED)
        add_arrow(s, x + 1.14, y + 0.41, 6.62, 3.09, RGBColor(0xC7, 0xD1, 0xCC), 0.7)

    # 05
    s = add_slide(prs, 5, "从入口到决策的页面结构", "首页承接高频执行，看板承接状态流转，统计与洞察承接低频复盘。", "INFORMATION ARCHITECTURE")
    add_panel(s, 1.05, 1.72, 11.25, 4.55, fill=PANEL)
    add_panel(s, 4.95, 2.0, 3.4, 0.56, fill=PALE_GREEN, line=RGBColor(0xC7, 0xE6, 0xDA))
    add_text(s, "智能 TodoLife", 5.58, 2.18, 2.1, 0.14, 12, True, GREEN, PP_ALIGN.CENTER)
    tabs = [("今日", GREEN), ("看板", BLUE), ("日历", ORANGE), ("任务", RED), ("我的", GRAY)]
    for i, (t, c) in enumerate(tabs):
        x = 1.65 + i * 2.05
        add_panel(s, x, 3.05, 1.35, 0.56, fill=RGBColor(0xF8, 0xFA, 0xF9), line=RGBColor(0xDE, 0xE7, 0xE2))
        add_text(s, t, x, 3.22, 1.35, 0.12, 10, True, c, PP_ALIGN.CENTER)
        add_arrow(s, 6.65, 2.56, x + 0.68, 3.05, RGBColor(0xC7, 0xD1, 0xCC), 0.8)
    overlays = [("AI 助手", BLUE), ("新建任务", GREEN), ("高级洞察", RED), ("统计页", ORANGE), ("倒计时页", GRAY)]
    for i, (t, c) in enumerate(overlays):
        x = 1.65 + i * 2.05
        add_label(s, t, x, 4.62, 1.35, 0.36, c, RGBColor(0xF6, 0xF8, 0xF7), 8.3)
        add_arrow(s, x + 0.68, 3.61, x + 0.68, 4.62, RGBColor(0xC7, 0xD1, 0xCC), 0.8)
    add_text(s, "高频执行", 1.8, 5.55, 2.3, 0.16, 9, True, GREEN)
    add_text(s, "状态流转", 4.02, 5.55, 2.3, 0.16, 9, True, BLUE)
    add_text(s, "时间管理", 6.05, 5.55, 2.3, 0.16, 9, True, ORANGE)
    add_text(s, "低频复盘", 8.16, 5.55, 2.3, 0.16, 9, True, RED)

    # 06
    s = add_slide(prs, 6, "任务不只是标题", "SmartTask 把一个待办扩展为可排序、可拆解、可提醒、可复盘的行为单元。", "DATA MODEL")
    groups = [
        ("基础信息", "标题、描述、分类、标签", GREEN),
        ("执行信息", "预计时长、截止日期、提醒时间、子任务", BLUE),
        ("决策信息", "重要度、复杂度、优先级、目标时间", ORANGE),
        ("状态信息", "Todo / In Progress / Done、完成时间", RED),
        ("习惯信息", "是否习惯、连续天数、完成历史", GRAY),
    ]
    for i, (t, b, c) in enumerate(groups):
        add_bullet(s, t, b, 0.85, 1.55 + i * 0.86, 4.2, 0.64, c, PANEL)
    add_panel(s, 7.0, 1.42, 4.45, 4.45, fill=PANEL)
    add_text(s, "SmartTask 卡片剖面", 7.36, 1.72, 2.5, 0.22, 14, True, TEXT)
    add_task_bar(s, 7.36, 2.22, 3.55, 0.48, "完成期末论文初稿", GREEN, 10)
    rows = [("category", "STUDY"), ("estimatedMinutes", "90"), ("dueDate", "2026-06-20"), ("priority", "HIGH"), ("status", "IN_PROGRESS"), ("subtasks", "3 / 5")]
    for i, (k, v) in enumerate(rows):
        y = 2.98 + i * 0.39
        add_text(s, k, 7.42, y, 1.45, 0.13, 7.6, False, MUTED)
        add_text(s, v, 9.35, y, 1.45, 0.13, 7.8, True, TEXT, PP_ALIGN.RIGHT)
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(7.42), inch(y + 0.22), inch(10.85), inch(y + 0.22))
        ln.line.color.rgb = RGBColor(0xEA, 0xEE, 0xEC)
        ln.line.width = Pt(0.4)

    # 07
    s = add_slide(prs, 7, "从一句话变成可执行计划", "远程大模型优先，本地启发式方案兜底；AI 失败时核心拆解仍然可用。", "AI PLANNING")
    flow = [
        ("用户输入", "下周五前完成论文初稿", GREEN),
        ("DueDateParser", "解析截止日期", BLUE),
        ("SmartTaskEngine", "推断分类 / 复杂度 / 时长", ORANGE),
        ("AI Planner", "拆解子任务", RED),
    ]
    for i, (t, b, c) in enumerate(flow):
        x = 0.75 + i * 2.75
        add_panel(s, x, 2.12, 2.2, 1.0, fill=PANEL)
        add_text(s, t, x + 0.18, 2.34, 1.75, 0.15, 9.5, True, c)
        add_text(s, b, x + 0.18, 2.64, 1.78, 0.22, 7.8, False, MUTED)
        if i < len(flow) - 1:
            add_arrow(s, x + 2.22, 2.62, x + 2.72, 2.62, RGBColor(0xB9, 0xC6, 0xC0), 1.2)
    add_panel(s, 2.0, 4.35, 9.25, 1.32, fill=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, "生成带计划日期的任务清单", 2.35, 4.6, 2.7, 0.2, 12, True, TEXT)
    tasks = [("资料调研", "6/10", "45 分钟"), ("梳理大纲", "6/12", "40 分钟"), ("撰写初稿", "6/15", "90 分钟"), ("修改完善", "6/18", "60 分钟"), ("排版检查", "6/20", "30 分钟")]
    for i, (name, date, mins) in enumerate(tasks):
        add_task_bar(s, 2.42 + i * 1.68, 5.02, 1.32, 0.34, name, [GREEN, BLUE, ORANGE, RED, GRAY][i], 7.2)
        add_text(s, f"{date} · {mins}", 2.44 + i * 1.68, 5.42, 1.28, 0.12, 5.9, False, MUTED, PP_ALIGN.CENTER)

    # 08
    s = add_slide(prs, 8, "让任务安排适配人的状态", "认知负荷不是装饰性指标，它决定界面是否需要收敛，以及任务是否适合现在推进。", "COGNITIVE MODEL")
    dims = [
        ("视觉负荷", "活跃任务数量、分类分散度", GREEN),
        ("记忆负荷", "高复杂任务数量、当前专注度", BLUE),
        ("时间压力", "临近目标时间任务、压力水平", ORANGE),
        ("决策疲劳", "任务总量、压力、精力状态", RED),
    ]
    for i, (t, b, c) in enumerate(dims):
        x = 0.9 + (i % 2) * 3.4
        y = 1.78 + (i // 2) * 1.28
        add_bullet(s, t, b, x, y, 2.85, 0.88, c, PANEL)
    cx, cy = 9.48, 3.08
    levels = [("低负荷", GREEN), ("平衡", RGBColor(0xD3, 0xA7, 0x00)), ("高负荷", ORANGE), ("过载", RED)]
    for i, (lab, col) in enumerate(levels):
        add_label(s, lab, 8.0 + i * 0.88, 1.62, 0.7, 0.28, col, RGBColor(0xF8, 0xFA, 0xF9), 7)
    radar_pts = [(cx, cy - 1.22), (cx + 1.22, cy), (cx, cy + 1.22), (cx - 1.22, cy)]
    for rr in [0.35, 0.7, 1.05]:
        poly = s.shapes.add_shape(MSO_SHAPE.DIAMOND, inch(cx - rr), inch(cy - rr), inch(2 * rr), inch(2 * rr))
        poly.fill.background()
        set_line(poly, RGBColor(0xD9, 0xE2, 0xDE), 0.7, 15)
    for x, y in radar_pts:
        add_arrow(s, cx, cy, x, y, RGBColor(0xD9, 0xE2, 0xDE), 0.6)
    vals = [(cx, cy - 0.82), (cx + 0.95, cy), (cx, cy + 0.55), (cx - 0.68, cy)]
    min_x = min(p[0] for p in vals)
    min_y = min(p[1] for p in vals)
    scale = inch(1) / 1000
    free = s.shapes.build_freeform(
        int((vals[0][0] - min_x) * 1000),
        int((vals[0][1] - min_y) * 1000),
        scale=scale,
    )
    free.add_line_segments(
        [(int((x - min_x) * 1000), int((y - min_y) * 1000)) for x, y in vals[1:]],
        close=True,
    )
    poly = free.convert_to_shape(inch(min_x), inch(min_y))
    set_fill(poly, RGBColor(0xA7, 0xD7, 0xCA), 35)
    set_line(poly, GREEN, 1.3)
    add_panel(s, 1.05, 5.55, 10.9, 0.6, fill=PALE_GREEN, line=RGBColor(0xC7, 0xE6, 0xDA))
    add_text(s, "当系统判断负荷偏高时，会切换到更简洁的专注视图，减少用户继续筛选任务的成本。", 1.38, 5.77, 10.1, 0.14, 10.5, True, GREEN, PP_ALIGN.CENTER)

    # 09
    s = add_slide(prs, 9, "下一步推荐不是随机排序", "推荐分数同时考虑任务本身和用户状态，避免在高压力、低精力时强推高复杂任务。", "PRIORITY LOGIC")
    factors = [("重要度", 36, GREEN, "任务价值"), ("紧急度", 30, ORANGE, "目标时间"), ("能力匹配", 20, BLUE, "专注与精力"), ("复杂度", 14, RED, "执行难度"), ("习惯加成", 8, GRAY, "连续性")]
    x0, y0, total_w = 0.95, 2.02, 6.1
    acc = 0
    for lab, pct, col, note in factors:
        w = total_w * pct / 108
        seg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x0 + acc), inch(y0), inch(w), inch(0.42))
        set_fill(seg, col)
        seg.line.fill.background()
        add_text(s, f"{pct}%", x0 + acc, y0 + 0.14, w, 0.11, 7.3, True, RGBColor(0xFF, 0xFF, 0xFF), PP_ALIGN.CENTER)
        acc += w
    for i, (lab, pct, col, note) in enumerate(factors):
        add_factor_card(s, 0.95 + (i % 3) * 2.1, 3.02 + (i // 3) * 1.0, lab, f"{pct}%", col, note)
    add_panel(s, 8.0, 2.0, 3.5, 2.85, fill=PANEL)
    add_text(s, "现在做", 8.35, 2.34, 0.8, 0.16, 9, True, GREEN)
    add_text(s, "整理项目汇报材料", 8.35, 2.68, 2.5, 0.22, 14, True, TEXT)
    add_label(s, "42 分钟", 8.35, 3.14, 0.88, 0.3, GREEN, PALE_GREEN, 8)
    add_label(s, "推荐分 82.4", 9.38, 3.14, 1.05, 0.3, BLUE, PALE_BLUE, 8)
    add_text(s, "主要因为重要度高、距离目标时间较近，同时当前专注与精力能够承载中等复杂任务。", 8.35, 3.68, 2.72, 0.54, 9.2, False, MUTED)

    # 10
    s = add_slide(prs, 10, "把任务池变成行动策略", "系统不只是分类任务，还把每一类转译成可执行策略。", "EISENHOWER MATRIX")
    add_text(s, "重要程度", 0.82, 3.45, 0.8, 0.2, 9, True, MUTED)
    add_text(s, "紧急程度", 6.45, 6.18, 1.0, 0.2, 9, True, MUTED, PP_ALIGN.CENTER)
    quads = [
        (1.72, 1.65, "重要且紧急", "立即做", "优先清空，避免扩散为风险", RED, PALE_RED,
         ["今晚前提交论文修改", "客户方案最终确认"]),
        (6.58, 1.65, "重要不紧急", "计划做", "投入收益最高，安排专注时段", GREEN, PALE_GREEN,
         ["完成期末论文初稿", "复盘本周学习计划"]),
        (1.72, 4.04, "紧急不重要", "限时做 / 委托", "快速处理，不占用深度精力", ORANGE, PALE_ORANGE,
         ["回复群内确认消息", "整理临时会议纪要"]),
        (6.58, 4.04, "不重要不紧急", "删除 / 批处理", "减少噪音，降低任务池负担", GRAY, PALE_GRAY,
         ["清理低优先级收藏", "批量整理截图"]),
    ]
    for x, y, title, action, desc, col, fill, tasks in quads:
        add_panel(s, x, y, 4.25, 1.92, fill=fill, line=col)
        add_text(s, title, x + 0.24, y + 0.18, 1.55, 0.18, 11, True, col)
        add_label(s, action, x + 2.75, y + 0.16, 1.05, 0.28, col, PANEL, 7.8)
        add_text(s, desc, x + 0.24, y + 0.52, 3.35, 0.18, 7.6, False, MUTED)
        for i, task in enumerate(tasks):
            add_task_bar(s, x + 0.27, y + 0.88 + i * 0.42, 3.55, 0.28, task, col, 7)
            add_text(s, f"{30 + i * 15} 分钟", x + 3.07, y + 0.96 + i * 0.42, 0.62, 0.1, 5.8, False, MUTED, PP_ALIGN.RIGHT)

    # 11
    s = add_slide(prs, 11, "首页只回答一个问题：现在做什么", "首页不是信息越多越好，而是把用户从“选择任务”带到“开始行动”。", "TODAY EXPERIENCE")
    add_phone(s, 2.1, 1.35, 2.55, 5.35, "今日", "today")
    notes = [
        ("顶部活力环", "完成率 / 专注分 / 习惯连续", 5.55, 1.55, GREEN),
        ("当前推荐", "把最该做的一件事压缩成高亮行动条", 5.95, 2.58, BLUE),
        ("快速捕捉", "单行输入，输入后展开备注与解析标签", 5.65, 3.62, ORANGE),
        ("状态与排程", "认知负荷、状态调节、智能时间段", 5.25, 4.72, RED),
    ]
    for title, body, x, y, col in notes:
        add_bullet(s, title, body, x, y, 4.6, 0.72, col, PANEL)
        add_arrow(s, x - 0.15, y + 0.36, 4.78, y + 0.15, col, 0.9)

    # 12
    s = add_slide(prs, 12, "不同阶段，用不同视图", "今日、看板、日历与倒计时分别解决执行、推进和时间管理问题。", "MULTI-VIEW")
    views = [("今日视图", "面向执行，按逾期 / 今天 / 稍后组织", "today", GREEN),
             ("看板视图", "面向推进，待处理 / 进行中 / 已完成", "board", BLUE),
             ("日历与倒计时", "面向时间，管理截止日期和重要节点", "calendar", ORANGE)]
    for i, (title, desc, mode, col) in enumerate(views):
        x = 1.05 + i * 4.05
        add_phone(s, x + 0.56, 1.52, 2.05, 4.35, title[:2], mode)
        add_text(s, title, x, 6.05, 3.18, 0.18, 12, True, col, PP_ALIGN.CENTER)
        add_text(s, desc, x, 6.35, 3.18, 0.23, 8.4, False, MUTED, PP_ALIGN.CENTER)

    # 13
    s = add_slide(prs, 13, "复盘不是报表，是下一轮安排的依据", "统计页帮助用户发现负荷高峰、低效时间段和稳定习惯，而不是展示复杂 BI 大屏。", "ANALYTICS")
    add_panel(s, 0.82, 1.55, 11.72, 4.88, fill=PANEL)
    add_text(s, "周报摘要", 1.15, 1.86, 1.4, 0.18, 11, True, TEXT)
    for i, (v, lab, col) in enumerate([("68%", "完成率", GREEN), ("17/25", "已完成", BLUE), ("上升", "趋势", ORANGE)]):
        add_panel(s, 1.15 + i * 1.35, 2.25, 1.06, 0.72, fill=RGBColor(0xF7, 0xFA, 0xF8))
        add_text(s, v, 1.15 + i * 1.35, 2.42, 1.06, 0.2, 13, True, col, PP_ALIGN.CENTER)
        add_text(s, lab, 1.15 + i * 1.35, 2.69, 1.06, 0.12, 7, False, MUTED, PP_ALIGN.CENTER)
    add_text(s, "完成率趋势", 1.15, 3.35, 1.4, 0.16, 10, True, TEXT)
    add_mini_line_chart(s, 1.15, 3.68, 3.8, 1.15, GREEN)
    add_text(s, "时间分布", 5.38, 1.86, 1.2, 0.18, 10, True, TEXT)
    colors = [GREEN, BLUE, RED, RGBColor(0x8B, 0x5C, 0xF6), ORANGE, RGBColor(0xEC, 0x48, 0x99)]
    labels = ["工作", "学习", "健康", "生活", "协作", "创作"]
    for i, (lab, col) in enumerate(zip(labels, colors)):
        add_label(s, lab, 5.38 + (i % 2) * 1.05, 2.25 + (i // 2) * 0.45, 0.82, 0.3, col, RGBColor(0xF8, 0xFA, 0xF9), 7.2)
    add_text(s, "认知负荷曲线", 8.0, 1.86, 1.6, 0.18, 10, True, TEXT)
    add_mini_line_chart(s, 8.0, 2.25, 3.25, 1.15, BLUE)
    add_text(s, "习惯热力图", 5.38, 4.0, 1.4, 0.18, 10, True, TEXT)
    add_heatmap(s, 5.4, 4.38, 15, 4, 0.13)
    add_text(s, "成就徽章", 8.0, 4.0, 1.4, 0.18, 10, True, TEXT)
    for i, (name, col) in enumerate([("十项全能", GREEN), ("一周习惯", ORANGE), ("低负荷管理", BLUE)]):
        add_label(s, name, 8.0, 4.36 + i * 0.48, 1.62, 0.3, col, RGBColor(0xF8, 0xFA, 0xF9), 7)

    # 14
    s = add_slide(prs, 14, "端侧为主，AI 增强", "架构重点是本地稳定可用 + AI 能力增强；AI 失败时，本地规则仍能支撑核心任务拆解和推荐。", "TECH ARCHITECTURE")
    layers = [
        ("UI 层", "Jetpack Compose / Material 3 / Screens / Components", GREEN),
        ("状态层", "ViewModel / StateFlow", BLUE),
        ("领域层", "SmartTaskEngine / CognitiveLoadEngine / StatisticsEngine / AchievementEngine / TaskCompletionEngine", ORANGE),
        ("数据层", "Local Repository / Theme Preference / Countdown Repository", RED),
        ("AI 层", "Remote AI Planner / Local AI Planner / AI Chat Client / DueDateParser", BLUE),
        ("系统能力", "ReminderScheduler / Notification / BroadcastReceiver", GRAY),
    ]
    for i, (name, body, col) in enumerate(layers):
        y = 1.45 + i * 0.78
        add_panel(s, 1.0, y, 11.0, 0.52, fill=PANEL)
        add_label(s, name, 1.25, y + 0.12, 1.08, 0.26, col, RGBColor(0xF8, 0xFA, 0xF9), 7.5)
        add_text(s, body, 2.6, y + 0.16, 8.7, 0.12, 8.5, False, TEXT)
        if i < len(layers) - 1:
            add_arrow(s, 6.5, y + 0.52, 6.5, y + 0.78, RGBColor(0xC7, 0xD1, 0xCC), 0.8)
    add_panel(s, 1.35, 6.2, 10.25, 0.45, fill=PALE_GREEN, line=RGBColor(0xC7, 0xE6, 0xDA))
    add_text(s, "AI 能力作为增强层接入，不成为基础任务管理的单点依赖。", 1.7, 6.34, 9.45, 0.12, 9.4, True, GREEN, PP_ALIGN.CENTER)

    # 15
    s = add_slide(prs, 15, "智能 TodoLife 的设计价值", "真正的效率工具，不是让用户管理更多事项，而是帮助用户在合适的时间做合适的事。", "DESIGN SUMMARY")
    values = [
        ("01", "从任务记录转向行动决策", GREEN),
        ("02", "从固定清单转向状态自适应", BLUE),
        ("03", "从单点 AI 转向可降级的智能流程", ORANGE),
        ("04", "从完成任务转向长期行为洞察", RED),
    ]
    for i, (num, text, col) in enumerate(values):
        y = 1.68 + i * 0.92
        add_panel(s, 1.42, y, 10.05, 0.66, fill=PANEL)
        add_text(s, num, 1.75, y + 0.19, 0.5, 0.14, 10.5, True, col)
        add_text(s, text, 2.48, y + 0.17, 6.8, 0.18, 14, True, TEXT)
    add_panel(s, 2.2, 5.78, 8.9, 0.66, fill=PALE_GREEN, line=RGBColor(0xC7, 0xE6, 0xDA))
    add_text(s, "把任务转译为行动策略，用数据反馈修正下一轮安排。", 2.55, 6.0, 8.2, 0.15, 12, True, GREEN, PP_ALIGN.CENTER)

    prs.save(OUT)


if __name__ == "__main__":
    build_ppt()
